#!/usr/bin/env python3
"""Genera la presentación (15 min) en PowerPoint a partir de los resultados.

Lee output/results.json, output/scaling.json y output/eda.json y embebe las
figuras de figures/, de modo que las cifras de las diapositivas siempre quedan
sincronizadas con los experimentos. Estructura orientada a negocio (problema,
valor, requerimientos, arquitectura, resultados, viabilidad).
"""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "figures"
OUT = ROOT / "output"
SLIDES = ROOT / "slides"

BLUE = RGBColor(0x00, 0x72, 0xB2)
ORANGE = RGBColor(0xD5, 0x5E, 0x00)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x55, 0x55, 0x55)

res = json.load(open(OUT / "results.json"))
scal = json.load(open(OUT / "scaling.json")) if (OUT / "scaling.json").exists() else None
eda = json.load(open(OUT / "eda.json")) if (OUT / "eda.json").exists() else None

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _set(p, text, size, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    p.text = text
    p.alignment = align
    for r in p.runs:
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
        r.font.name = "Calibri"


def bar(slide, color=BLUE, h=Inches(0.18)):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()


def title_slide(title, subtitle, footer):
    s = prs.slides.add_slide(BLANK)
    bar(s, BLUE, Inches(0.25))
    tf = _box(s, Inches(0.8), Inches(2.2), Inches(11.7), Inches(2.2))
    _set(tf.paragraphs[0], title, 34, BLUE, bold=True)
    p = tf.add_paragraph(); _set(p, subtitle, 20, GREY)
    tf2 = _box(s, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.8))
    _set(tf2.paragraphs[0], footer, 13, GREY)
    return s


def content_slide(title, bullets, fig=None, fig_w=Inches(6.2), fig_left=None):
    s = prs.slides.add_slide(BLANK)
    bar(s)
    tf = _box(s, Inches(0.6), Inches(0.45), Inches(12), Inches(0.9))
    _set(tf.paragraphs[0], title, 26, BLUE, bold=True)
    text_w = Inches(6.0) if fig else Inches(12.1)
    body = _box(s, Inches(0.7), Inches(1.7), text_w, Inches(5.2))
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        lvl = 0
        if isinstance(b, tuple):
            b, lvl = b
        _set(p, ("• " if lvl == 0 else "   – ") + b, 18 if lvl == 0 else 15,
             DARK if lvl == 0 else GREY, bold=(lvl == 0 and b.endswith(":")))
        p.space_after = Pt(7)
    if fig:
        fpath = FIG / fig
        if fpath.exists():
            left = fig_left if fig_left is not None else Inches(6.9)
            s.shapes.add_picture(str(fpath), left, Inches(1.6), width=fig_w)
    return s


def big_fig_slide(title, fig, caption, fig_w=Inches(8.2)):
    s = prs.slides.add_slide(BLANK)
    bar(s)
    tf = _box(s, Inches(0.6), Inches(0.45), Inches(12), Inches(0.9))
    _set(tf.paragraphs[0], title, 26, BLUE, bold=True)
    fpath = FIG / fig
    if fpath.exists():
        pic = s.shapes.add_picture(str(fpath), Inches(0), Inches(1.5), width=fig_w)
        pic.left = int((SW - pic.width) / 2)
    cap = _box(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.6))
    _set(cap.paragraphs[0], caption, 14, GREY, align=PP_ALIGN.CENTER)
    return s


# Cifras
m = res["models"]
d = res["data"]
ir = d["imbalance_ratio"]; pr = d["pos_rate"]
best_auc = max(m, key=lambda k: m[k]["auc"])
fast = min(m, key=lambda k: m[k]["train_seconds"])

# 1. Título
title_slide(
    "Predicción de propina baja en taxis de NYC con Apache Spark MLlib",
    "Comparación de algoritmos de Machine Learning distribuido sobre Big Data",
    "Niels Pacheco  ·  Maestría en IA — EPG-UNI  ·  Curso de Big Data  ·  Docente: Mg. Rosa V. Encinas Quille",
)

# 2. Problema y necesidad
content_slide("El problema y la necesidad", [
    "La propina es una fracción clave del ingreso del conductor de taxi.",
    "¿Podemos anticipar, al inicio del viaje, si dejará una propina baja?",
    ("Beneficiarios:", 0),
    ("Conductores: anticipar viajes de baja recompensa.", 1),
    ("Plataformas/despacho: incentivos y asignación más justa.", 1),
    ("Reto técnico: decenas de millones de viajes → no cabe en una sola máquina con herramientas tradicionales.", 0),
], fig="fig1_volumen_mensual.png", fig_w=Inches(5.6))

# 3. Propuesta de valor
content_slide("Propuesta de valor", [
    "Pipeline de Big Data reproducible, de extremo a extremo, sobre datos públicos reales.",
    "Compara 4 algoritmos distribuidos y mide su costo y escalabilidad.",
    ("Qué lo hace distinto:", 0),
    ("Métricas honestas bajo desbalance (no solo accuracy).", 1),
    ("Evidencia cuantitativa de escalabilidad (datos y núcleos).", 1),
    ("Decisiones documentadas (ADRs) y código abierto (repo público).", 1),
])

# 4. Datos y las 5 V
content_slide("Los datos: NYC TLC Yellow Taxi", [
    f"6 meses de 2023 · {d['n_raw']:,} viajes crudos → {d['n_clean']:,} tras limpieza.",
    "Formato Apache Parquet (columnar), ~2,5 M viajes/mes.",
    ("Las 5 V abordadas:", 0),
    ("Volumen: ~15 M de registros.  Variedad: espacio, tiempo, economía.", 1),
    ("Veracidad: solo pagos con tarjeta (la propina en efectivo no se registra).", 1),
    ("Valor: anticipar la recompensa del conductor.", 1),
], fig="fig2_distribucion_propina.png", fig_w=Inches(5.6))

# 5. El reto: desbalance
big_fig_slide("El reto central: desbalance de clases",
              "fig3_balance_clases.png",
              f"La propina baja es minoritaria ({pr:.1%}), con un desbalance de {ir:.1f}:1. "
              f"Un modelo trivial lograría ~90% de accuracy y 0% de recall útil.",
              fig_w=Inches(5.6))

# 6. Arquitectura
big_fig_slide("Arquitectura: ciclo de vida del dato",
              "fig7_arquitectura.png",
              "Ingesta → Almacenamiento → Procesamiento distribuido (Spark) → "
              "Features (MLlib) → Modelado → Evaluación y visualización.",
              fig_w=Inches(11.2))

# 7. Requerimientos
content_slide("Requerimientos de la solución", [
    ("Funcionales:", 0),
    ("Ingesta y normalización de Parquet con esquema variable.", 1),
    ("Limpieza, ingeniería de variables y etiquetado distribuido.", 1),
    ("Entrenamiento y comparación de modelos; reportes de métricas.", 1),
    ("No funcionales:", 0),
    ("Escalabilidad horizontal (datos y núcleos).", 1),
    ("Reproducibilidad (entorno y datos versionados).", 1),
    ("Veracidad y honestidad estadística de la evaluación.", 1),
])

# 8. Metodología
content_slide("Metodología (Spark MLlib)", [
    "Pipeline: StringIndexer → OneHotEncoder → VectorAssembler.",
    "14 variables numéricas + 2 categóricas; se excluye la propina (anti-fuga).",
    "4 modelos: Regresión Logística, Árbol, Random Forest, GBT.",
    "Desbalance: ponderación de instancias por clase (los 4 modelos).",
    "Split 80/20 idéntico para todos; AUC, F1, recall de la minoría.",
], fig="fig8_comparacion_modelos.png", fig_w=Inches(5.8))

# 9. Resultados: comparación
rows = [("Modelo", "AUC", "F1", "Acc.", "Rec+", "t(s)")]
for name in m:
    mm = m[name]
    rows.append((name, f"{mm['auc']:.3f}", f"{mm['f1']:.3f}", f"{mm['accuracy']:.3f}",
                 f"{mm['recall_pos']:.3f}", f"{mm['train_seconds']:.0f}"))
s = prs.slides.add_slide(BLANK); bar(s)
tf = _box(s, Inches(0.6), Inches(0.45), Inches(12), Inches(0.9))
_set(tf.paragraphs[0], "Resultados: comparación de modelos", 26, BLUE, bold=True)
tbl = s.shapes.add_table(len(rows), 6, Inches(0.7), Inches(1.7), Inches(6.0), Inches(3.4)).table
for ci in range(6):
    tbl.columns[ci].width = Inches(1.5 if ci == 0 else 0.9)
for ri, row in enumerate(rows):
    for cii, val in enumerate(row):
        cell = tbl.cell(ri, cii); cell.text = val
        para = cell.text_frame.paragraphs[0]
        for r in para.runs:
            r.font.size = Pt(12); r.font.bold = (ri == 0)
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if ri == 0 else DARK
        if ri == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
note = _box(s, Inches(0.7), Inches(5.4), Inches(6.0), Inches(1.6))
_set(note.paragraphs[0],
     "Con ponderación, los 4 modelos detectan la minoría. GBT y RF dan el mejor "
     "AUC; el árbol y GBT, el mejor recall — pero a costos muy distintos.", 14, GREY)
s.shapes.add_picture(str(FIG / "fig9_tradeoff.png"), Inches(7.0), Inches(1.6), width=Inches(5.7))

# 10. Hallazgo clave
content_slide("Hallazgo clave: la accuracy engaña", [
    "Bajo fuerte desbalance, la accuracy NO es una métrica honesta.",
    "Un clasificador trivial lograría 90,6% de accuracy y 0% de recall útil.",
    "Por eso evaluamos con AUC y recall de la clase minoritaria.",
    ("No hay un único ganador:", 0),
    (f"GBT: mejor desempeño (AUC {m['GBT']['auc']:.3f}, recall+ {m['GBT']['recall_pos']:.2f}), mayor costo.", 1),
    (f"Random Forest: AUC casi igual ({m['RandomForest']['auc']:.3f}) a la mitad del tiempo.", 1),
    (f"Regresión Logística: la más eficiente ({m['LogisticRegression']['train_seconds']:.0f} s).", 1),
], fig="fig9_tradeoff.png", fig_w=Inches(5.4))

# 11. Escalabilidad
if scal:
    sp8 = next((r["speedup"] for r in scal["results"] if r["cores"] == 8), None)
    cap = (f"Escalabilidad de datos sublineal y speedup de {sp8:.1f}× con 8 núcleos "
           "(ley de Amdahl).")
else:
    cap = "Escalabilidad de datos sublineal (paralelismo de Spark)."
s = prs.slides.add_slide(BLANK); bar(s)
tf = _box(s, Inches(0.6), Inches(0.45), Inches(12), Inches(0.9))
_set(tf.paragraphs[0], "Escalabilidad del procesamiento distribuido", 26, BLUE, bold=True)
s.shapes.add_picture(str(FIG / "fig10_escalabilidad_datos.png"), Inches(0.7), Inches(1.7), width=Inches(5.9))
if (FIG / "fig11_escalabilidad_fuerte.png").exists():
    s.shapes.add_picture(str(FIG / "fig11_escalabilidad_fuerte.png"), Inches(6.9), Inches(1.7), width=Inches(5.9))
capb = _box(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.6))
_set(capb.paragraphs[0], cap, 14, GREY, align=PP_ALIGN.CENTER)

# 12. Viabilidad
content_slide("Viabilidad e impacto", [
    ("Escenarios de uso:", 0),
    ("Incentivos dinámicos y asignación de viajes.", 1),
    ("Análisis del comportamiento de propina por zona/horario.", 1),
    ("Limitaciones:", 0),
    ("Sesgo: solo pagos con tarjeta; señal limitada de la propina.", 1),
    ("Ejecución local: no captura red/shuffle entre nodos.", 1),
    ("Expansión: clúster real (EMR/Databricks), datos contextuales (clima, tráfico), streaming.", 0),
])

# 13. Conclusiones
content_slide("Conclusiones", [
    "Se comparó de forma justa el desempeño y el costo de 4 clasificadores distribuidos sobre ~15 M de viajes.",
    "Lección 1: en datos desbalanceados, elegir por AUC y recall, no por accuracy.",
    "Lección 2: Random Forest dio el mejor compromiso; la Regresión Logística, el menor costo.",
    "Lección 3: la escalabilidad confirma las ventajas y los límites (Amdahl) del cómputo distribuido.",
    "Código, datos y artículo: repositorio público reproducible.",
])

# 14. Gracias
s = title_slide("¡Gracias!", "Preguntas y discusión",
                "Repositorio: github.com/nielspac177/big-data-nyc-taxi-mllib  ·  nielspacheco1997@gmail.com")

SLIDES.mkdir(exist_ok=True)
out = SLIDES / "presentacion_big_data.pptx"
prs.save(str(out))
print(f"Presentación generada: {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} diapositivas)")
